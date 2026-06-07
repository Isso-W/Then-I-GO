# -*- coding: utf-8 -*-
"""Generate PPT matching the 12-slide white+purple template."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── palette (white bg, purple accent) ──
BG_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE    = RGBColor(0x6C, 0x5C, 0xFF)
PURPLE_LT = RGBColor(0xA9, 0x8B, 0xFF)
PURPLE_BG = RGBColor(0xF3, 0xF0, 0xFF)   # light purple card bg
CARD_BG   = RGBColor(0xF7, 0xF7, 0xFA)   # light gray card bg
BLACK     = RGBColor(0x1A, 0x1A, 0x2E)
GRAY      = RGBColor(0x55, 0x55, 0x66)
LIGHT_G   = RGBColor(0x99, 0x99, 0xAA)
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
CYAN      = RGBColor(0x06, 0xB6, 0xD4)
RED       = RGBColor(0xEF, 0x44, 0x44)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GOLD      = RGBColor(0xFF, 0xD1, 0x66)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG_WHITE):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, radius=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def line(slide, l, t, w, color=PURPLE, thickness=Pt(3)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, thickness)
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    return shp

def tb(slide, l, t, w, h, text, sz=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return txBox

def bullets(slide, items, l, t, w, h, sz=15, color=GRAY, sp=Pt(8)):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color; p.space_after = sp
    return txBox

def number_badge(slide, num, l, t, sz=Inches(0.55)):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    shp.fill.solid(); shp.fill.fore_color.rgb = PURPLE; shp.line.fill.background()
    tf = shp.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(20); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0); tf.paragraphs[0].space_after = Pt(0)
    shp.text_frame.margin_top = Pt(4)

def slide_header(slide, num, title, subtitle=""):
    number_badge(slide, num, Inches(0.6), Inches(0.45))
    tb(slide, Inches(1.3), Inches(0.4), Inches(8), Inches(0.6), title, sz=32, color=BLACK, bold=True)
    line(slide, Inches(1.3), Inches(1.05), Inches(1.8))
    if subtitle:
        tb(slide, Inches(1.3), Inches(1.15), Inches(10), Inches(0.45), subtitle, sz=14, color=LIGHT_G)

def card(slide, l, t, w, h, title, body, title_color=PURPLE, bg=CARD_BG, body_color=GRAY, body_sz=13):
    rect(slide, l, t, w, h, bg)
    tb(slide, l + Inches(0.25), t + Inches(0.15), w - Inches(0.5), Inches(0.45),
       title, sz=15, color=title_color, bold=True)
    tb(slide, l + Inches(0.25), t + Inches(0.55), w - Inches(0.5), h - Inches(0.7),
       body, sz=body_sz, color=body_color)

def card_purple(slide, l, t, w, h, title, body):
    card(slide, l, t, w, h, title, body, title_color=PURPLE, bg=PURPLE_BG)

# ════════════════════════════════════════
# 1. 技术架构总览
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 1, "技术架构总览", "那我走 Then-I-GO  |  Architecture Overview")

# 中心流程：用户 → 画像 → 评分 → Agent → 路线 → 探索 → 反馈 → 回环
flow_labels = ["用户输入", "画像模块", "POI 评分", "AI Agent", "路线生成", "探索体验", "行为采集", "自学习回环"]
flow_colors = [PURPLE_BG, PURPLE_BG, PURPLE_BG, PURPLE_BG, PURPLE_BG, PURPLE_BG, PURPLE_BG, PURPLE_BG]
start_x = Inches(0.6)
card_w = Inches(1.35)
gap = Inches(0.17)
y = Inches(2.8)
for i, label in enumerate(flow_labels):
    x = start_x + i * (card_w + gap)
    rect(s, x, y, card_w, Inches(0.7), PURPLE_BG)
    tb(s, x, y + Inches(0.12), card_w, Inches(0.5), label, sz=13, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(flow_labels) - 1:
        tb(s, x + card_w, y + Inches(0.1), Inches(0.17), Inches(0.5), ">", sz=16, color=PURPLE_LT, align=PP_ALIGN.CENTER)

# 上层标注
tb(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.7),
   "纯前端架构  |  Vite 插件代理 Gemini API  |  零后端服务器  |  localStorage 持久化  |  React + Tailwind + motion/react",
   sz=15, color=GRAY)

# 下层详细模块
modules = [
    ("Onboarding\nMBTI 收集", "冷启动单步"),
    ("12 维评分\n确定性过滤", "零 LLM 调用"),
    ("路线 Agent\n审查 Agent\n对话 Agent", "3 个 AI Agent"),
    ("ReAct 循环\n生成-审查-修正", "最多 1 轮"),
    ("动态站点\n隐藏任务\nA/B 岔路", "intensity 分档"),
    ("Vlog 生成\n战报卡\n足迹地图", "零视频模型"),
    ("TripRecord\nemoji 反馈", "行为信号"),
    ("偏好注入\n(预埋)", "闭环就绪"),
]
for i, (top_text, bottom_text) in enumerate(modules):
    x = start_x + i * (card_w + gap)
    tb(s, x, Inches(3.8), card_w, Inches(1.5), top_text, sz=11, color=GRAY, align=PP_ALIGN.CENTER)
    tb(s, x, Inches(5.3), card_w, Inches(0.4), bottom_text, sz=10, color=LIGHT_G, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════
# 2. 用户画像模块
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 2, "用户画像模块", "User Understanding Engine")

card(s, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.3),
     "冷启动 Onboarding", (
         "首次启动仅收集 MBTI（单步，16 型四列网格）\n"
         "存入 localStorage，无需注册/登录\n\n"
         "口味偏好步骤已移除 — 降低进入门槛\n"
         "路线生成时 MBTI 作为长期画像注入 prompt"
     ))

card(s, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.3),
     "5 步偏好收集（每次探索前）", (
         "1. 心情 — 开心/疲惫/无聊/放松/探索/想吃（6选1）\n"
         "2. 时长 — 1h / 2h / 半天 / 整天\n"
         "3. 兴趣标签 — 文艺/户外/美食/热闹/拍照/小众/省钱\n"
         "4. 同行人 — 独自/情侣/朋友/亲子\n"
         "5. 安排程度 — 别让我思考(默认) / 正常探索"
     ))

card_purple(s, Inches(0.6), Inches(4.5), Inches(3.7), Inches(2.2),
     "「别让我思考」模式", (
         "隐藏站名+描述，到达才揭晓\n"
         "无 A/B 岔路选择，系统直选\n"
         "强避重（去过的 -0.8）\n\n"
         "产品洞察：大多数用户说「那我走」\n"
         "的时候，就是不想做任何决定"
     ))

card(s, Inches(4.6), Inches(4.5), Inches(3.7), Inches(2.2),
     "「正常探索」模式", (
         "全部显示站名、描述、奖励\n"
         "A/B 岔路选择（气质相反的两个候选）\n"
         "弱避重（去过的 -0.4）\n\n"
         "想要掌控感的用户主动切换"
     ))

card(s, Inches(8.6), Inches(4.5), Inches(3.7), Inches(2.2),
     "设置页 MBTI 修改", (
         "设置页可随时修改 MBTI\n"
         "16 种 MBTI 4x4 网格选择\n"
         "选中即保存到 localStorage\n\n"
         "影响后续所有路线生成的 AI prompt"
     ))

# ════════════════════════════════════════
# 3. AI Agent 决策层
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 3, "AI Agent 决策层", "AI Agent Route Planning")

tb(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
   "3 个 AI Agent（调用 Gemini，具备推理与行动能力）+ 2 个确定性辅助模块（纯规则计算）",
   sz=16, color=GRAY)

# 三个 Agent 卡片
card_purple(s, Inches(0.6), Inches(2.3), Inches(3.7), Inches(2.5),
     "路线 Agent (routeAgent)", (
         "从候选 POI 选点、排序、编排站点序列\n"
         "生成故事文案 / 打卡任务 / 奖励\n"
         "支持 revision：带审查意见重新生成\n\n"
         "模型：gemini-2.5-flash-lite\n"
         "输出：严格 JSON，正则提取"
     ))

card_purple(s, Inches(4.6), Inches(2.3), Inches(3.7), Inches(2.5),
     "审查 Agent (routeReviewer)", (
         "品类多样性 — 连续同类型？\n"
         "活动节奏 — 饭后紧接剧烈运动？\n"
         "时间合理性 — 上午排酒吧？\n"
         "体验递进 — 节奏是否单调？\n"
         "综合常识 — 连逛三家小卖铺？"
     ))

card_purple(s, Inches(8.6), Inches(2.3), Inches(3.7), Inches(2.5),
     "对话 Agent (chatAgent)", (
         "探索中自然语言聊天改路线\n"
         "4 种操作指令：\n"
         "  replace_next — 替换下一站\n"
         "  skip_current — 跳过当前站\n"
         "  add_stop — 末尾加站\n"
         "  none — 纯聊天\n"
         "候选来自真实 POI 池，不会编造"
     ))

# 辅助模块
card(s, Inches(0.6), Inches(5.2), Inches(5.8), Inches(1.6),
     "POI 评分过滤 (poiFilter) — 确定性模块", (
         "12 维评分函数 · 零 LLM 调用 · 渐进降级过滤\n"
         "每个品类最多 2 个（category cap）· 输出 top 15 候选给路线 Agent"
     ), title_color=CYAN)

card(s, Inches(6.8), Inches(5.2), Inches(5.8), Inches(1.6),
     "用户画像 — 数据透传", (
         "MBTI + 兴趣标签存 localStorage · 路线生成时注入 prompt\n"
         "不是独立 Agent，是上下文数据层"
     ), title_color=CYAN)

# ════════════════════════════════════════
# 4. 探索路线生成
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 4, "探索路线生成", "Dynamic Route Generation  |  ReAct Loop")

# 流程图
steps = [
    ("1. POI 过滤", "poiFilter\n12 维评分\ntop 15 候选", CARD_BG, PURPLE),
    ("2. 路线生成", "routeAgent\nGemini 选点排序\n故事/任务/奖励", PURPLE_BG, PURPLE),
    ("3. ReAct 审查", "routeReviewer\n品类/节奏/时间\n通过 or 修正", PURPLE_BG, PURPLE),
    ("4. 营业校验", "确定性预检\n步行+停留时间\n打烊站 LLM 换补", CARD_BG, PURPLE),
    ("5. 输出路线", "waypoints[]\nhiddenTask\nbranch?", CARD_BG, GREEN),
]

for i, (title, body, bg, tc) in enumerate(steps):
    x = Inches(0.4) + i * Inches(2.55)
    rect(s, x, Inches(1.8), Inches(2.3), Inches(2.5), bg)
    tb(s, x, Inches(1.9), Inches(2.3), Inches(0.4), title, sz=14, color=tc, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x, Inches(2.4), Inches(2.3), Inches(1.8), body, sz=12, color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        tb(s, x + Inches(2.3), Inches(2.6), Inches(0.25), Inches(0.5), ">", sz=20, color=PURPLE_LT, align=PP_ALIGN.CENTER)

# 下半部分 - 关键特性
card(s, Inches(0.6), Inches(4.7), Inches(3.7), Inches(2.1),
     "动态站点数量", (
         "1 小时 → 1 主线站 + 1 隐藏任务\n"
         "2 小时 → 1 主线站 + 1 隐藏任务\n"
         "半天   → 2 主线站 + 1 隐藏任务\n"
         "整天   → 3 主线站 + 1 隐藏任务"
     ))

card(s, Inches(4.6), Inches(4.7), Inches(3.7), Inches(2.1),
     "隐藏任务生成策略", (
         "Gemini 一次调用多生成 1 个站点\n"
         "最后一个 pop() 出来作为隐藏任务\n"
         "坐标真实 / 品类不冲突 / 可跳过\n"
         "奖励要求更好、任务更有挑战"
     ))

card(s, Inches(8.6), Inches(4.7), Inches(3.7), Inches(2.1),
     "A/B 岔路抉择", (
         "第一站打卡后弹出两个气质相反的候选\n"
         "（安静书店 vs 热闹酒吧）\n"
         "附一句抉择提示文案\n"
         "仅「正常探索」模式启用"
     ))

# ════════════════════════════════════════
# 5. 探索服务系统
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 5, "探索服务系统", "Exploration Space System")

card(s, Inches(0.6), Inches(1.8), Inches(3.7), Inches(2.3),
     "地图与导航", (
         "SVG 地图：terrain 底色 + 真实路网\n"
         "Dijkstra 沿街最短路寻路\n"
         "紫色虚线 + 方向箭头导航\n"
         "双端点选路避免回头路\n"
         "拖动/点击地图 → 吸附最近路 → 走过去"
     ))

card(s, Inches(4.6), Inches(1.8), Inches(3.7), Inches(2.3),
     "方向面板", (
         "旋转箭头图标实时指向目标\n"
         "距离 + 预计到达时间\n"
         "不用文字方向（上午/东/西无意义）\n"
         "Heading 追踪：移动 >2m 更新方向\n"
         "箭头角度 = 行走方向到目标的相对夹角"
     ))

card(s, Inches(8.6), Inches(1.8), Inches(3.7), Inches(2.3),
     "接近打卡（模拟 LBS）", (
         "打卡按钮不拦截，随时可拍照记录\n"
         "30m 内亮绿色提示「到了 · 可打卡」\n"
         "否则显示「距目标 Nm」\n\n"
         "给到达感，但不强制"
     ))

card(s, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.5),
     "探索状态机", (
         "intro  >  preference_selection  >  gear_confirmation\n"
         "  > waypointIndex=0: initial > checkin_initial\n"
         "    > hidden_found > hidden_active > checkin_hidden > reward_hidden\n"
         "    > branch_choice > onAdvanceWaypoint(+1)\n"
         "  > waypointIndex=1..N: next_objective > checkin_next > advance\n"
         "  > achievement_unlock > feedback > intro\n\n"
         "位置由 positionFromStep() 推导，不瞬移"
     ))

card(s, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.5),
     "聊天改路线", (
         "右下角紫色气泡 > 底部半屏对话面板\n"
         "AI 感知：当前路线上下文 + 附近 8 个未用候选 POI\n"
         "4 种操作：replace_next / skip_current / add_stop / none\n"
         "操作结果以系统消息插入聊天记录\n\n"
         "候选来自 poiFilter 实时过滤，确保推荐真实存在"
     ))

# ════════════════════════════════════════
# 6. 奖励系统
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 6, "奖励系统", "Gamification Engine")

card_purple(s, Inches(0.6), Inches(1.8), Inches(3.7), Inches(2.2),
     "优惠券奖励", (
         "每站打卡奖励与品类相关的优惠券\n"
         "咖啡厅 > ¥15 饮品抵扣券\n"
         "书店 > ¥20 购书折扣券\n"
         "统一「¥金额+券名」格式\n"
         "不发徽章、不发 XP — 只发有价值的东西"
     ))

card_purple(s, Inches(4.6), Inches(1.8), Inches(3.7), Inches(2.2),
     "品类锚定 Gemini 输出", (
         "CATEGORY_REWARDS 品类>奖励映射表\n"
         "prompt 给每个候选 POI 附推荐奖励\n"
         "含具体金额参考\n"
         "Gemini 可微调文案但有锚点\n"
         "避免「恭喜获得探索勋章」空洞文案"
     ))

card_purple(s, Inches(8.6), Inches(1.8), Inches(3.7), Inches(2.2),
     "成就徽章系统", (
         "8 枚成就徽章，基于真实 tripHistory\n"
         "连续打卡 / 隐藏任务达人\n"
         "集齐品类 / 集券数量\n"
         "背包页展示 + 优惠券 QR 二维码弹窗"
     ))

# 品类图片
card(s, Inches(0.6), Inches(4.4), Inches(11.7), Inches(2.5),
     "18 张品类图（public/categories/）", (
         "餐饮 9 类：地方菜系 / 火锅 / 烧烤 / 异国料理 / 自助餐 / 鱼鲜 / 小吃 / 饮品 / 甜品\n"
         "非餐饮 9 类：咖啡 / 酒吧 / 书店 / 文创 / 公园 / 健身 / 娱乐 / 骑行券 / 打车券\n\n"
         "getCategoryImage(category) 把 POI category 映射到对应图片路径\n"
         "用于背包优惠券卡片和隐藏任务图片"
     ))

# ════════════════════════════════════════
# 7. AI Vlog 生成
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 7, "AI Vlog 生成", "AI Vlog Generation  |  Zero Video Model")

tb(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.5),
   "拆解「视频」为各自便宜的部分 — per-user x per-second 成本 demo 扛不住，所以不调视频生成模型",
   sz=15, color=GRAY)

# 4 个组件
comps = [
    ("分镜表", "路线 waypoints\n直接当分镜", "免费"),
    ("路径几何", "Dijkstra 寻路\n折线插值", "免费"),
    ("旁白/字幕/BGM", "一次 Gemini\n文本补全", "~$0.001"),
    ("画面运动", "地图回放 +\nLive2D 呼吸", "0 token"),
]
for i, (title, body, cost) in enumerate(comps):
    x = Inches(0.6) + i * Inches(3.1)
    card_purple(s, x, Inches(2.2), Inches(2.8), Inches(1.8), title, body + "\n\n" + cost)

# 下半
card(s, Inches(0.6), Inches(4.4), Inches(3.7), Inches(2.5),
     "三种风格", (
         "赛博朋克\n"
         "  霓虹扫描线 + Synthwave 110BPM\n"
         "复古胶片\n"
         "  噪点暗角 + City Pop 90BPM\n"
         "日系清新\n"
         "  漏光效果 + Lo-fi 75BPM"
     ))

card(s, Inches(4.6), Inches(4.4), Inches(3.7), Inches(2.5),
     "回放体验", (
         "1. 地图回放\n"
         "   小人沿路网走 + 轨迹按风格色生长\n"
         "   到站标记弹出 + IG Stories 进度条\n"
         "2. 到站特写\n"
         "   Live2D 呼吸动画 + 字幕旁白\n"
         "3. 战报卡\n"
         "   标题+金句+路线缩略图+数据格子"
     ))

card(s, Inches(8.6), Inches(4.4), Inches(3.7), Inches(2.5),
     "后续升级方向", (
         "Vlog Agent 完整化\n"
         "  接入 image-to-video 模型\n"
         "  生成真实视频片段\n"
         "  多轮迭代优化剪辑节奏\n\n"
         "用户上传素材替换占位 B-roll\n"
         "  目前 3 张预生成占位图"
     ), title_color=AMBER)

# ════════════════════════════════════════
# 8. 数据回环
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 8, "数据回环", "Behavioral Feedback Loop  |  RLHF")

card_purple(s, Inches(0.6), Inches(1.8), Inches(3.7), Inches(2.3),
     "行为信号采集", (
         "每站是否实际到访 (visited)\n"
         "A/B 岔路选了哪个\n"
         "聊天改路线操作日志\n"
         "是否触发/跳过隐藏任务\n"
         "探索总时长和距离\n"
         "用户 emoji 一键评价"
     ))

card_purple(s, Inches(4.6), Inches(1.8), Inches(3.7), Inches(2.3),
     "极轻反馈设计", (
         "「跟我吐槽一下」浮层\n"
         "一行 emoji（6 个表情）\n"
         "点一个不到 1 秒就完成\n"
         "5 秒不操作自动消失\n\n"
         "历史页可给每站单独打 emoji\n"
         "反馈门槛低到不存在"
     ))

card_purple(s, Inches(8.6), Inches(1.8), Inches(3.7), Inches(2.3),
     "足迹地图", (
         "个人页展示所有到访站点\n"
         "SVG mask 迷雾效果\n"
         "去过的点「挖洞」透出底图\n"
         "未探索区域被迷雾覆盖\n"
         "探索越多 地图越清晰\n"
         "适配明暗主题 (--fog-color)"
     ))

# RLHF
card(s, Inches(0.6), Inches(4.5), Inches(11.7), Inches(2.3),
     "RLHF — Reinforcement Learning from Human Feedback", (
         "用户的行为即反馈，系统无需额外询问就能做出一致性调整\n\n"
         "示例：用户跳过隐藏任务 >\n"
         "  故事页素材集 — 不显示该站点      |  Vlog 配图列表 — 移除该站\n"
         "  历史记录 — 条目消失                    |  结算奖励 — 不计入\n\n"
         "不是跳过了还到处看到「你错过了一个隐藏任务」 — 跳过就是跳过，全链路尊重用户决定"
     ), title_color=PURPLE, bg=PURPLE_BG)

# ════════════════════════════════════════
# 9. 系统架构
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 9, "系统架构", "Frontend Architecture  |  Data Layer")

card(s, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.2),
     "前端架构", (
         "React + Tailwind CSS v4 + motion/react\n"
         "Screen routing: App.tsx 内 screen state，无路由库\n"
         "6 个主屏幕：explore / story / bag / mine / event / settings\n"
         "App-level state: 路线/偏好/探索状态机/主题/历史/Vlog"
     ))

card(s, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.2),
     "Gemini 调用链路", (
         "前端 agent > src/lib/gemini.ts > POST /api/generate\n"
         "Vite 插件 vite-plugin-gemini-proxy.ts 代理\n"
         "优先 Vertex AI (GOOGLE_APPLICATION_CREDENTIALS)\n"
         "回退 AI Studio (GEMINI_API_KEY)"
     ))

card(s, Inches(0.6), Inches(4.4), Inches(3.7), Inches(2.5),
     "地理数据层", (
         "路网: 34 条路段, 28.2km\n"
         "  scripts/data/street-network.json\n"
         "地表: 11 块多边形\n"
         "  scripts/data/terrain.json\n"
         "POI: 200 条, 31 品类\n"
         "  src/data/pois.json\n"
         "全部 commit 进仓库，完全离线"
     ))

card(s, Inches(4.6), Inches(4.4), Inches(3.7), Inches(2.5),
     "持久化 (localStorage)", (
         "userProfile — MBTI + 兴趣标签\n"
         "tripHistory — TripRecord[]\n"
         "vlogHistory — GeneratedVlog[]\n"
         "confirmedGear — 装备清单\n"
         "appTheme — 主题偏好\n\n"
         "预埋 3 条样本 TripRecord"
     ))

card(s, Inches(8.6), Inches(4.4), Inches(3.7), Inches(2.5),
     "主题系统", (
         "深色/浅色主题切换\n"
         "首次按系统时间自动选\n"
         "  6:00-18:00 日间，其余夜间\n"
         "CSS 变量 + data-theme 属性\n"
         "手动切换后 localStorage 持久化\n"
         "Vlog 地图/路网/地表色跟随"
     ))

# ════════════════════════════════════════
# 10. 技术亮点
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 10, "技术亮点", "Technical Highlights")

highlights = [
    ("纯前端，无后端", "Vite 插件代理 Gemini，生产环境零服务器"),
    ("ReAct 审查循环", "AI 生成 > AI 审查 > 带反馈重生成"),
    ("确定性时间预检", "审查前先算步行+停留，超预算 50% 直接打回"),
    ("12 维评分系统", "mood x weather x time x distance x novelty"),
    ("RLHF 行为响应反馈", "跳过隐藏任务 > 全链路同步移除"),
    ("全链路可观测日志", "DevTools Console 追踪完整 AI 决策链路"),
    ("200 POI 人工审查", "Gemini 生成 + 逐条 tag/mood/price 校验"),
    ("路网 Dijkstra + snap", "所有 POI 吸附路面，不穿墙不飞线"),
    ("Vlog 零视频模型", "文本补全 + 程序化动画，成本趋近于零"),
    ("自然语言改路线", "聊天 Agent 只从真实 POI 池选，不编造"),
]

for i, (title, desc) in enumerate(highlights):
    col = Inches(0.6) if i % 2 == 0 else Inches(6.8)
    row_y = Inches(1.7) + (i // 2) * Inches(1.1)
    rect(s, col, row_y, Inches(5.8), Inches(0.95), PURPLE_BG if i % 2 == 0 else CARD_BG)
    tb(s, col + Inches(0.3), row_y + Inches(0.1), Inches(5.2), Inches(0.4),
       title, sz=16, color=PURPLE, bold=True)
    tb(s, col + Inches(0.3), row_y + Inches(0.5), Inches(5.2), Inches(0.4),
       desc, sz=13, color=GRAY)

# ════════════════════════════════════════
# 11. 与传统推荐平台对比
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 11, "与传统推荐平台对比", "Why Different")

# 表头
header_y = Inches(1.8)
rect(s, Inches(0.6), header_y, Inches(3.2), Inches(0.6), PURPLE)
tb(s, Inches(0.6), header_y + Inches(0.1), Inches(3.2), Inches(0.4), "对比维度", sz=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
rect(s, Inches(3.8), header_y, Inches(4.2), Inches(0.6), CARD_BG)
tb(s, Inches(3.8), header_y + Inches(0.1), Inches(4.2), Inches(0.4), "传统推荐平台", sz=15, color=GRAY, bold=True, align=PP_ALIGN.CENTER)
rect(s, Inches(8.0), header_y, Inches(4.6), Inches(0.6), PURPLE_BG)
tb(s, Inches(8.0), header_y + Inches(0.1), Inches(4.6), Inches(0.4), "那我走", sz=15, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

rows = [
    ("决策模式", "用户自己挑选、比较、做攻略", "说「那我走」，AI 帮你决定一切"),
    ("推荐逻辑", "评分排序 / 协同过滤", "12 维评分 x 多 Agent 协作生成"),
    ("路线编排", "无（散点 POI 列表）", "站点序列 + 故事文案 + 打卡任务"),
    ("实时调整", "无", "自然语言聊天改路线（4 种操作）"),
    ("个性化深度", "口味 / 价格", "MBTI + 心情 + 天气 + 时段 + 同行人"),
    ("惊喜感", "无", "隐藏任务 + A/B 岔路 + mystery 模式"),
    ("行为闭环", "评分 / 收藏", "RLHF 全链路响应 + 极轻 emoji 反馈"),
    ("内容产出", "用户写点评", "AI 自动生成 Vlog + 战报卡"),
]

for i, (dim, trad, ours) in enumerate(rows):
    y = Inches(2.45) + i * Inches(0.6)
    bg = CARD_BG if i % 2 == 0 else BG_WHITE
    rect(s, Inches(0.6), y, Inches(3.2), Inches(0.55), bg)
    tb(s, Inches(0.8), y + Inches(0.08), Inches(2.8), Inches(0.4), dim, sz=13, color=BLACK, bold=True)
    rect(s, Inches(3.8), y, Inches(4.2), Inches(0.55), bg)
    tb(s, Inches(4.0), y + Inches(0.08), Inches(3.8), Inches(0.4), trad, sz=12, color=LIGHT_G)
    rect(s, Inches(8.0), y, Inches(4.6), Inches(0.55), PURPLE_BG if i % 2 == 0 else BG_WHITE)
    tb(s, Inches(8.2), y + Inches(0.08), Inches(4.2), Inches(0.4), ours, sz=12, color=PURPLE)

# ════════════════════════════════════════
# 12. 未来规划
# ════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, 12, "未来规划", "What's Next")

future = [
    ("Vlog Agent 完整化", "接入 image-to-video 模型，生成真实视频片段，多轮迭代优化剪辑节奏"),
    ("用户上传素材", "用户拍摄/上传真实照片替换占位 B-roll，Vlog 画面真实化"),
    ("真实 GPS LBS", "替换模拟 LBS，接入 navigator.geolocation 实现真实位置打卡"),
    ("偏好自学习接入", "buildHistorySummary 注入 prompt，路线推荐越来越个性化"),
    ("社交分享奖励", "分享 Vlog 到外部平台获得 in-app 奖励回写"),
    ("自然语言 feedback", "文本框直接吐槽，AI 理解反馈调整偏好权重"),
]

for i, (title, desc) in enumerate(future):
    col = Inches(0.6) if i % 2 == 0 else Inches(6.8)
    row_y = Inches(1.7) + (i // 2) * Inches(1.5)
    card_purple(s, col, row_y, Inches(5.8), Inches(1.3), title, desc)

# 结语
tb(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.7),
   "走出去，发现城市的下一面。", sz=24, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)

# ── save ──
out = r"D:\Then-I-GO\docs\那我走-项目汇报v2.pptx"
prs.save(out)
print("PPT done:", out)
